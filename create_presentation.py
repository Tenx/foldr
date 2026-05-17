#!/usr/bin/env python3
"""
Create a professional PowerPoint presentation for Frontend Slides project
Using Creative Bold style (Google/Airbnb inspired)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Creative Bold Brand Colors
PRIMARY = RGBColor(255, 107, 53)  # Vibrant Orange #FF6B35
SECONDARY = RGBColor(41, 128, 185)  # Ocean Blue #2980B9
DARK = RGBColor(44, 62, 80)  # Dark Blue-Gray #2C3E50
LIGHT = RGBColor(236, 240, 241)  # Off-White #ECF0F1
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(33, 33, 33)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Add vibrant accent shape
    accent_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0), Inches(0), Inches(6), Inches(9)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = PRIMARY
    accent_shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(6.5), Inches(3), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Frontend Slides"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(84)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK
    title_para.font.name = "Montserrat"

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(6.5), Inches(5.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Create Stunning HTML Presentations\nwith Zero Dependencies"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = DARK
    subtitle_para.font.name = "Montserrat"
    subtitle_para.line_spacing = 1.3

    # Slide 2: The Problem
    slide = add_content_slide(prs, "The Problem",
        "Most presentation tools force you to choose:",
        [
            "Design quality OR ease of use",
            "Dependencies OR modern features",
            "Templates OR customization",
            "Describe your vision OR see what you get"
        ], SECONDARY)

    # Slide 3: The Solution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK

    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Frontend Slides: Show, Don't Tell"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(64)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.font.name = "Montserrat"
    title_para.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(12), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "A Claude Code skill that creates beautiful,\nanimation-rich HTML presentations through visual exploration"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(36)
    subtitle_para.font.color.rgb = LIGHT
    subtitle_para.font.name = "Montserrat"
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.line_spacing = 1.4

    # Slide 4: Core Philosophy
    slide = add_four_point_slide(prs, "Core Philosophy",
        [
            ("🎯 Zero Dependencies", "Single HTML files. No npm, no build tools."),
            ("👁️ Show, Don't Tell", "Pick from visual previews."),
            ("✨ Anti-AI-Slop", "12 curated distinctive styles."),
            ("📱 Production Quality", "Accessible, responsive, customizable.")
        ], PRIMARY)

    # Slide 5: How It Works
    slide = add_content_slide(prs, "How It Works", "",
        [
            '1. Tell Claude: "Create a pitch deck"',
            "2. Describe your content and feeling",
            "3. See 3 visual style previews",
            "4. Get your presentation, auto-opened"
        ], SECONDARY)

    # Slide 6: 12 Curated Styles - Dark
    slide = add_two_column_slide(prs, "12 Curated Styles",
        "Dark Themes",
        [
            "Bold Signal - High-impact, vibrant",
            "Electric Studio - Professional, clean",
            "Creative Voltage - Energetic, neon",
            "Dark Botanical - Elegant, sophisticated"
        ],
        "Light Themes",
        [
            "Notebook Tabs - Editorial, organized",
            "Pastel Geometry - Friendly, approachable",
            "Split Pastel - Playful, modern",
            "Vintage Editorial - Witty, personality"
        ], PRIMARY, SECONDARY)

    # Slide 7: PowerPoint Conversion
    slide = add_content_slide(prs, "PowerPoint Conversion",
        "Transform PPT to Web in 4 Steps",
        [
            "1. Extract - All text, images, and notes",
            "2. Preview - Confirm extracted content",
            "3. Style - Pick your visual aesthetic",
            "4. Generate - HTML with all assets"
        ], PRIMARY)

    # Slide 8: Key Features
    slide = add_four_point_slide(prs, "Key Features",
        [
            ("📐 Viewport Fitting", "Every slide fits in 100vh. No scrolling."),
            ("🎨 Distinctive Design", "Curated fonts, bold colors, atmosphere."),
            ("🎬 Rich Animations", "CSS-only micro-interactions."),
            ("📝 Well-Commented", "Future-you will thank present-you.")
        ], SECONDARY)

    # Slide 9: Use Cases
    slide = add_two_column_slide(prs, "Use Cases",
        "Perfect For:",
        [
            "🚀 Startup pitch decks",
            "📊 Product launches",
            "🎓 Teaching and tutorials",
            "💼 Conference talks"
        ],
        "Why Choose HTML?",
        [
            "Works offline forever",
            "Easy to embed and share",
            "Lightweight (50-200KB)",
            "Full control over design"
        ], PRIMARY, SECONDARY)

    # Slide 10: Installation
    slide = add_big_quote_slide(prs,
        "/plugin marketplace add zarazhangrui/frontend-slides\n/plugin install frontend-slides@frontend-slides",
        "Installation - Two Commands",
        DARK)

    # Slide 11: Architecture
    slide = add_content_slide(prs, "Progressive Disclosure Architecture",
        "Main skill: ~180 lines (concise workflow map)",
        [
            "Supporting files loaded on-demand:",
            "• Style presets (12 themes)",
            "• Viewport responsive CSS",
            "• Animation patterns",
            "• PPT extraction script"
        ], SECONDARY)

    # Slide 12: Share Your Work
    slide = add_two_column_slide(prs, "Share Your Work",
        "Deploy to Live URL",
        [
            "bash scripts/deploy.sh",
            "→ Permanent Vercel URL",
            "→ Works on all devices"
        ],
        "Export to PDF",
        [
            "bash scripts/export-pdf.sh",
            "→ For email, Slack, Notion",
            "→ Static snapshot"
        ], PRIMARY, SECONDARY)

    # Slide 13: Design Philosophy
    slide = add_philosophy_slide(prs, [
        ("You don't need to be a designer", "You just need to react to what you see"),
        ("Dependencies are debt", "A single HTML file works in 10 years"),
        ("Generic is forgettable", "Every presentation feels custom-crafted")
    ], PRIMARY)

    # Slide 14: Why Frontend Slides?
    slide = add_comparison_slide(prs, "Why Frontend Slides?",
        [
            ("❌ Generic AI Templates", "✅ Curated Distinctive Design"),
            ("Purple gradients, Inter font", "Bold choices, unique fonts"),
            ("❌ Describe What You Want", "✅ Show and Pick"),
            ('"Make it modern..."', "Here are 3 options. Which one?")
        ], PRIMARY)

    # Slide 15: Get Started
    slide = add_content_slide(prs, "Get Started Today",
        "",
        [
            "1. Install: /plugin marketplace add zarazhangrui/frontend-slides",
            "2. Run: /frontend-slides",
            "3. Pick a style you love",
            "4. Share with the world 🚀"
        ], PRIMARY)

    # Slide 16: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PRIMARY

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Thank You"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(96)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.font.name = "Montserrat"
    title_para.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(12), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "github.com/zarazhangrui/frontend-slides\n\nLet's make the web beautiful ✨"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.font.name = "Montserrat"
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.line_spacing = 1.5

    return prs

def add_content_slide(prs, title, subtitle, bullets, accent_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.3), Inches(9)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent_color
    accent_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK
    title_para.font.name = "Montserrat"

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(14), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(32)
        subtitle_para.font.color.rgb = accent_color
        subtitle_para.font.name = "Montserrat"

    # Bullets
    start_y = 3.2 if subtitle else 2.5
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(start_y), Inches(13), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i > 0:
            content_frame.add_paragraph()
        p = content_frame.paragraphs[i]
        p.text = bullet
        p.font.size = Pt(28)
        p.font.color.rgb = DARK
        p.font.name = "Montserrat"
        p.space_after = Pt(20)
        p.level = 0

    return slide

def add_four_point_slide(prs, title, points, accent_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK
    title_para.font.name = "Montserrat"
    title_para.alignment = PP_ALIGN.CENTER

    # Four points in 2x2 grid
    positions = [
        (1, 2.5), (8.5, 2.5),
        (1, 5.5), (8.5, 5.5)
    ]

    for i, (point_title, point_desc) in enumerate(points):
        x, y = positions[i]

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(6.5), Inches(2.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = accent_color
        card.line.width = Pt(3)

        # Point title
        title_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.3), Inches(6), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = point_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = accent_color
        title_para.font.name = "Montserrat"

        # Point description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 1.2), Inches(6), Inches(1))
        desc_frame = desc_box.text_frame
        desc_frame.text = point_desc
        desc_frame.word_wrap = True
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(20)
        desc_para.font.color.rgb = DARK
        desc_para.font.name = "Montserrat"

    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets, color1, color2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK
    title_para.font.name = "Montserrat"
    title_para.alignment = PP_ALIGN.CENTER

    # Left column
    left_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(2.2), Inches(7), Inches(5.5)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = color1
    left_card.line.fill.background()

    left_title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(6.2), Inches(0.8))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = left_title
    left_title_para = left_title_frame.paragraphs[0]
    left_title_para.font.size = Pt(36)
    left_title_para.font.bold = True
    left_title_para.font.color.rgb = WHITE
    left_title_para.font.name = "Montserrat"

    left_content_box = slide.shapes.add_textbox(Inches(1.2), Inches(3.5), Inches(6.2), Inches(4))
    left_content_frame = left_content_box.text_frame
    left_content_frame.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i > 0:
            left_content_frame.add_paragraph()
        p = left_content_frame.paragraphs[i]
        p.text = bullet
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.font.name = "Montserrat"
        p.space_after = Pt(15)

    # Right column
    right_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.2), Inches(2.2), Inches(7), Inches(5.5)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = color2
    right_card.line.fill.background()

    right_title_box = slide.shapes.add_textbox(Inches(8.6), Inches(2.5), Inches(6.2), Inches(0.8))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = right_title
    right_title_para = right_title_frame.paragraphs[0]
    right_title_para.font.size = Pt(36)
    right_title_para.font.bold = True
    right_title_para.font.color.rgb = WHITE
    right_title_para.font.name = "Montserrat"

    right_content_box = slide.shapes.add_textbox(Inches(8.6), Inches(3.5), Inches(6.2), Inches(4))
    right_content_frame = right_content_box.text_frame
    right_content_frame.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i > 0:
            right_content_frame.add_paragraph()
        p = right_content_frame.paragraphs[i]
        p.text = bullet
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.font.name = "Montserrat"
        p.space_after = Pt(15)

    return slide

def add_big_quote_slide(prs, quote, attribution, bg_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color

    # Quote box
    quote_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(12), Inches(3))
    quote_frame = quote_box.text_frame
    quote_frame.text = quote
    quote_frame.word_wrap = True
    quote_para = quote_frame.paragraphs[0]
    quote_para.font.size = Pt(36)
    quote_para.font.color.rgb = WHITE
    quote_para.font.name = "Courier New"
    quote_para.alignment = PP_ALIGN.CENTER
    quote_para.line_spacing = 1.4

    # Attribution
    attr_box = slide.shapes.add_textbox(Inches(2), Inches(6), Inches(12), Inches(1))
    attr_frame = attr_box.text_frame
    attr_frame.text = attribution
    attr_para = attr_frame.paragraphs[0]
    attr_para.font.size = Pt(28)
    attr_para.font.color.rgb = PRIMARY
    attr_para.font.name = "Montserrat"
    attr_para.alignment = PP_ALIGN.CENTER

    return slide

def add_philosophy_slide(prs, philosophies, accent_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Design Philosophy"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK
    title_para.font.name = "Montserrat"
    title_para.alignment = PP_ALIGN.CENTER

    y_positions = [2.5, 4.5, 6.5]

    for i, (principle, explanation) in enumerate(philosophies):
        y = y_positions[i]

        # Principle
        principle_box = slide.shapes.add_textbox(Inches(2), Inches(y), Inches(12), Inches(0.6))
        principle_frame = principle_box.text_frame
        principle_frame.text = f'"{principle}"'
        principle_para = principle_frame.paragraphs[0]
        principle_para.font.size = Pt(32)
        principle_para.font.bold = True
        principle_para.font.color.rgb = accent_color
        principle_para.font.name = "Montserrat"

        # Explanation
        exp_box = slide.shapes.add_textbox(Inches(2), Inches(y + 0.7), Inches(12), Inches(0.5))
        exp_frame = exp_box.text_frame
        exp_frame.text = explanation
        exp_para = exp_frame.paragraphs[0]
        exp_para.font.size = Pt(24)
        exp_para.font.color.rgb = DARK
        exp_para.font.name = "Montserrat"

    return slide

def add_comparison_slide(prs, title, comparisons, accent_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK
    title_para.font.name = "Montserrat"
    title_para.alignment = PP_ALIGN.CENTER

    y_start = 2.5
    row_height = 1.3

    for i, (left, right) in enumerate(comparisons):
        y = y_start + (i * row_height)

        # Left (bad)
        left_box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(7), Inches(1))
        left_frame = left_box.text_frame
        left_frame.text = left
        left_frame.word_wrap = True
        left_para = left_frame.paragraphs[0]
        left_para.font.size = Pt(24)
        left_para.font.color.rgb = DARK
        left_para.font.name = "Montserrat"

        # Right (good)
        right_box = slide.shapes.add_textbox(Inches(8.5), Inches(y), Inches(7), Inches(1))
        right_frame = right_box.text_frame
        right_frame.text = right
        right_frame.word_wrap = True
        right_para = right_frame.paragraphs[0]
        right_para.font.size = Pt(24)
        right_para.font.bold = True
        right_para.font.color.rgb = accent_color
        right_para.font.name = "Montserrat"

    return slide

if __name__ == "__main__":
    prs = create_presentation()
    output_file = "/Users/I742076/.claude/projects/foldr/Frontend-Slides-Presentation.pptx"
    prs.save(output_file)
    print(f"✅ Presentation created: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
